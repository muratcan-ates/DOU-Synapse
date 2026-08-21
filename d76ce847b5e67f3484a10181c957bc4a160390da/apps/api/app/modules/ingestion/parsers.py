"""Belge ayrıştırıcıları.

Tasarımın tek amacı vardır: **kaynak konumunu kaybetmemek.** Her metin bloğu hangi
sayfadan, slayttan veya kod satırından geldiğini taşır; çünkü sistemin merkezi vaadi
"her cevap sayfa numarasıyla gelir" (PLAN.md P0 #5). Konum bilgisi burada kaybolursa
zincirin geri kalanı onu geri getiremez.

PyMuPDF ve python-pptx bilinçli tercihtir: hızlı, az bağımlılıklı ve sayfa/slayt
sınırlarını doğal olarak korurlar. Karmaşık düzenli belgeler için Docling ikinci aşamada
yalnızca sorunlu dosyalara uygulanacak fallback'tir (ARCHITECTURE.md §1).
"""

from __future__ import annotations

import io
import itertools
import re
from dataclasses import dataclass, field

from app.core.errors import ValidationError
from app.models.core import ChunkContentType
from app.modules.ingestion.validation import CODE_EXTENSIONS


@dataclass(frozen=True, slots=True)
class ParsedBlock:
    """Kaynağı bilinen tek bir metin bloğu."""

    text: str
    content_type: ChunkContentType = ChunkContentType.TEXT
    page_number: int | None = None
    slide_number: int | None = None
    section_title: str | None = None


@dataclass(slots=True)
class ParsedDocument:
    blocks: list[ParsedBlock] = field(default_factory=list)
    page_count: int | None = None


def _clean(text: str) -> str:
    """Satır sonlarını normalize eder, tekrarlı boşlukları toplar."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def parse_pdf(content: bytes) -> ParsedDocument:
    import pymupdf

    try:
        document = pymupdf.open(stream=content, filetype="pdf")
    except Exception as exc:
        raise ValidationError("PDF dosyası açılamadı. Dosya bozuk olabilir.") from exc

    blocks: list[ParsedBlock] = []
    with document:
        if document.needs_pass:
            raise ValidationError("Parola korumalı PDF'ler desteklenmiyor.")

        # `pymupdf.Document` çalışma zamanında yinelenebilir ama tip bilgisi bunu
        # söylemiyor; açık `page_count` döngüsü hem mypy'ı geçiyor hem niyeti
        # daha net anlatıyor. Bu iki hata `mypy app` koşumunu tamamen durduruyordu,
        # yani tüm paket için tip denetimi fiilen kapalıydı.
        for index in range(1, document.page_count + 1):
            page = document.load_page(index - 1)
            text = _clean(page.get_text("text"))
            if text:
                blocks.append(ParsedBlock(text=text, page_number=index))
        page_count = document.page_count

    if not blocks:
        raise ValidationError(
            "PDF içinde metin bulunamadı. Taranmış (görüntü tabanlı) PDF'ler için "
            "önce metin tanıma uygulanmalıdır."
        )
    return ParsedDocument(blocks=blocks, page_count=page_count)


def parse_pptx(content: bytes) -> ParsedDocument:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise ValidationError("Sunum dosyası açılamadı. Dosya bozuk olabilir.") from exc

    blocks: list[ParsedBlock] = []
    slide_count = 0
    for index, slide in enumerate(presentation.slides, start=1):
        slide_count = index
        title: str | None = None
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = _clean(slide.shapes.title.text) or None

        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = _clean(shape.text_frame.text)
                if text and text != title:
                    parts.append(text)
            elif shape.has_table:
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows
                ]
                table = _clean("\n".join(rows))
                if table:
                    blocks.append(
                        ParsedBlock(
                            text=table,
                            content_type=ChunkContentType.TABLE,
                            slide_number=index,
                            section_title=title,
                        )
                    )

        body = _clean("\n".join(parts))
        if title or body:
            blocks.append(
                ParsedBlock(
                    text=_clean(f"{title}\n{body}" if title else body),
                    slide_number=index,
                    section_title=title,
                )
            )

    if not blocks:
        raise ValidationError("Sunum içinde metin bulunamadı.")
    return ParsedDocument(blocks=blocks, page_count=slide_count)


_MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")


def parse_markdown(content: bytes) -> ParsedDocument:
    """Başlık hiyerarşisini koruyarak bölümlere ayırır.

    Bölüm başlığı `section_title` olarak taşınır; kullanıcıya "sayfa" yerine bölüm adı
    gösterilebilmesi için gereklidir.
    """
    text = content.decode("utf-8", errors="replace")
    blocks: list[ParsedBlock] = []
    heading_stack: list[str] = []
    buffer: list[str] = []

    def flush() -> None:
        body = _clean("\n".join(buffer))
        buffer.clear()
        if body:
            blocks.append(ParsedBlock(text=body, section_title=" > ".join(heading_stack) or None))

    for line in text.splitlines():
        match = _MARKDOWN_HEADING.match(line.strip())
        if match:
            flush()
            level = len(match.group(1))
            heading_stack[level - 1 :] = [match.group(2).strip()]
        else:
            buffer.append(line)
    flush()

    if not blocks:
        raise ValidationError("Markdown dosyasında metin bulunamadı.")
    return ParsedDocument(blocks=blocks)


# Fonksiyon/sınıf/metot başlangıcı: kod, mantıksal sınırlarından bölünsün diye kullanılır.
_CODE_BOUNDARY = re.compile(
    r"^\s*(?:"
    r"(?:async\s+)?def\s+\w+|class\s+\w+|"  # Python
    r"(?:public|private|protected|static|final|\s)*[\w<>\[\]]+\s+\w+\s*\([^;]*\)\s*\{|"  # Java/C++
    r"(?:export\s+)?(?:async\s+)?function\s+\w+|"  # JS/TS
    r"(?:export\s+)?(?:default\s+)?class\s+\w+"
    r")"
)
_SYMBOL_NAME = re.compile(r"(?:def|class|function)\s+(\w+)|(\w+)\s*\(")


def parse_code(content: bytes, extension: str) -> ParsedDocument:
    """Kod dosyasını fonksiyon/sınıf sınırlarından böler.

    Satır aralığı `section_title` içinde taşınır; öğrenciye "şu dosyanın 25-47. satırları"
    diyebilmek, kod inceleme sorularında (code_trace / bug_hunt) kaynağın anlamlı
    gösterilmesi için gerekir.
    """
    text = content.decode("utf-8", errors="replace")
    lines = text.splitlines()
    if not any(line.strip() for line in lines):
        raise ValidationError("Kod dosyası boş görünüyor.")

    boundaries = [i for i, line in enumerate(lines) if _CODE_BOUNDARY.match(line)]
    if not boundaries or boundaries[0] != 0:
        boundaries.insert(0, 0)
    boundaries.append(len(lines))

    blocks: list[ParsedBlock] = []
    for start, end in itertools.pairwise(boundaries):
        segment = "\n".join(lines[start:end]).strip("\n")
        if not segment.strip():
            continue
        match = _SYMBOL_NAME.search(lines[start]) if start < len(lines) else None
        symbol = next((g for g in match.groups() if g), None) if match else None
        label = f"satır {start + 1}-{end}"
        blocks.append(
            ParsedBlock(
                text=segment,
                content_type=ChunkContentType.CODE,
                section_title=f"{symbol} ({label})" if symbol else label,
            )
        )
    return ParsedDocument(blocks=blocks)


def parse_text(content: bytes) -> ParsedDocument:
    text = _clean(content.decode("utf-8", errors="replace"))
    if not text:
        raise ValidationError("Dosyada metin bulunamadı.")
    return ParsedDocument(blocks=[ParsedBlock(text=text)])


def parse(content: bytes, extension: str) -> ParsedDocument:
    """Uzantıya göre uygun ayrıştırıcıyı seçer."""
    if extension == ".pdf":
        return parse_pdf(content)
    if extension == ".pptx":
        return parse_pptx(content)
    if extension == ".md":
        return parse_markdown(content)
    if extension in CODE_EXTENSIONS:
        return parse_code(content, extension)
    return parse_text(content)
