"""Testlerin ortak kurulum fabrikaları — tek kaynak.

## Neden bu dosya var

`tests/` altındaki modüller birbirinden yardımcı içe aktarıyordu: `test_exam_lock`
sohbet hattı sahtesini `test_chat_api`'den, sınav kurulumunu `test_exams`'ten,
parça üreticisini `test_socratic`'ten alıyordu. Böyle on üç bağ vardı ve her biri
GİZLİ bir bağlılıktı. Bedeli tek yönlü değildi:

* Kendi dosyasındaki bir fixture'ı düzelten kişi, uzaktaki üç dosyanın kurduğu
  senaryoyu da değiştiriyordu ve bunu ancak koşumda, ilgisiz görünen bir hatayla
  öğreniyordu.
* Ters yön daha sinsiydi: `test_chat_api`'nin sahtesini "sadeleştirmek",
  `test_exam_lock`'un kanıtladığı şeyi sessizce zayıflatabiliyordu — o dosyada
  hiçbir satır değişmediği için inceleme sırasında görünmezdi.

Aynı baskının ikinci yüzü kopyalamaydı: uzaktaki bir test dosyasına bağlanmak
istemeyen herkes aynı kurulumu yerinde yeniden yazdı. "Ders oluştur" yardımcısı
yedi dosyada ayrı ayrı duruyordu ve kopyalar ayrışmıştı — beşi `str`, ikisi
`UUID` döndürüyor, bir tanesi üyelik yanıtını `201` yerine `(200, 201)` sayıyordu.
Ayrışan kopya, ortak kurulumun en pahalı hâlidir: aynı adı taşıyan iki yardımcı
farklı şeyler yapınca, testleri okuyan kişi hangisini okuduğunu bilemez.

Bu dosya iki sorunu birden bitirir. Kurulum bir kez burada tanımlanır; testler
birbirine değil buraya bağlanır. Bir kurulumun değişmesi artık tek bir yerde
görünür ve kimi etkilediği içe aktarım listesinden okunabilir.

## Sınırlar

* **Burada pytest fixture'ı YOKTUR.** Fixture'lar `conftest.py`'nin işidir.
  Buradakiler düz fonksiyon ve sınıflardır; bir fixture zincirine girmeden de
  çağrılabilirler, dolayısıyla "bu testi koşmak için hangi fixture'lar gerekli"
  sorusu testin kendi imzasında okunur kalır.
* **Eşik üçtür** (Anayasa XI): üç ya da daha fazla yerde tekrarlanan kurulum
  buraya taşınır. İki yerde geçen bir yardımcıyı buraya almak, kazandığı
  tekrarsızlıktan çoğunu dolaylılık olarak geri verir.

  İkide kalan bir şey yalnız iki gerekçeyle buraya girer ve bugün beş girişi
  bu istisnaya dayanıyor: (a) taşımak bir çapraz test-modülü bağını kaldırıyorsa
  (`make_pptx`), (b) iki kopya zaten AYRIŞMIŞSA — o durumda "iki kopya" sayısı
  yanıltıcıdır, ortada tek bir davranış yoktur (`FakeRetriever`: biri çağrı
  sayıyordu, diğeri saymıyordu). `Pipeline` ve onun parçaları (`FakeGenerator`,
  `FakeCitationGuardrail`, `install_pipeline`) tek kullanıcılıdır ama hattın
  kendisi buradadır; parçalarını dışarıda bırakmak hattı yarım bırakırdı.
* **Sahtelerin davranışı burada DEĞİŞMEZ.** Ayrışan kopyaların en genel biçimi
  alındı; dar kullanımların ihtiyacı varsayılan argümanla karşılanıyor. Bir
  fabrikanın parametresi olan her alan, bugün en az bir çağrı yerinde gerçekten
  değişen bir alandır — spekülatif parametre yok. (Bu bir iddia değil ölçüt:
  `document_id` ve `slide_number` 48 çağrının hiçbirinde verilmediği için
  `make_chunk`'tan düşürüldü.)
"""

from __future__ import annotations

import asyncio
import io
import json
import time
from collections.abc import Sequence
from dataclasses import dataclass
from typing import TYPE_CHECKING, Any
from uuid import UUID, uuid4

from httpx import AsyncClient

# `sql_text` olarak alınıyor: bu modülün en çok çağrılan fabrikası olan
# `make_chunk`'ın parametresi de `text` ve iki adın çakışması, SQL sabitlerini
# okuyan kişiyi durduruyordu.
from sqlalchemy import text as sql_text
from sqlalchemy.ext.asyncio import AsyncEngine

from app.contracts import (
    AnswerStatus,
    ChatMode,
    Citation,
    GeneratedAnswer,
    GuardrailVerdict,
    RetrievedChunk,
    SocraticStage,
)
from app.models.assessment import QuestionType
from app.modules.ingestion.embedding import HashingEmbeddingProvider
from app.schemas.assessment import AnswerFormat

if TYPE_CHECKING:
    from tests.conftest import UserFactory

# ---------------------------------------------------------------------------
# Uç üzerinden kurulum
#
# Bu üç fabrika kasten HTTP'den geçer. Ders, üyelik ve konu satırlarını doğrudan
# SQL ile yazmak daha hızlı olurdu, ama o zaman testlerin kurduğu durum uçların
# ürettiği durumla ayrışabilirdi ve hiçbir test bunu göremezdi.
# ---------------------------------------------------------------------------


async def create_course(
    client: AsyncClient,
    headers: dict[str, str],
    code: str,
    *,
    title: str | None = None,
) -> UUID:
    """Ders oluşturur ve kimliğini döndürür; çağıran eğitmen olur.

    KARAR — dönüş tipi `UUID`'dir. Yedi kopyanın beşi `str`, ikisi `UUID`
    döndürüyordu ve bu ayrışma testlerde `UUID(course_id)` / `str(course_id)`
    çevrimlerini rastgele yerlere serpmişti. Kimliğin kendi tipiyle taşınması
    doğru olan: yol parametresine gömülürken `f"{course_id}"` zaten `str`
    üretir, veritabanına bind edilirken de `UUID` gerekir. `str` bekleyen çağrı
    yerleri kendi tarafında `str(...)` yazar — ters yön (her yerde `UUID(...)`)
    hem daha kalabalık hem de bozuk bir kimliği geç fark eder.

    `title` yalnız ders koduna bakan testler için türetilir; başlığın kendisini
    sınayan çağrı yerleri onu açıkça verir.
    """
    response = await client.post(
        "/courses",
        json={"code": code, "title": title if title is not None else f"{code} Dersi"},
        headers=headers,
    )
    assert response.status_code == 201, response.text
    return UUID(response.json()["id"])


async def enroll_student(
    client: AsyncClient,
    headers: dict[str, str],
    course_id: UUID | str,
    email: str,
) -> None:
    """Kayıtlı bir kullanıcıyı derse öğrenci olarak ekler.

    `headers` eğitmenin başlığı olmalıdır; öğrenci üye ekleyemez ve bu kural
    `test_courses` tarafından ayrıca sınanır.

    `201` bekleniyor: kopyaların biri `(200, 201)` sayıyordu, ama uç her zaman
    `201` döndürüyor (`test_courses` bunu doğrudan iddia ediyor). Gevşek kontrol
    bir şey kazandırmıyor, uç bir gün `200`'e dönerse sessizce kabul ederdi.

    `role` parametre DEĞİL: bugün on yedi çağrı yerinin hepsi `student` yazıyor.
    Eğitmen ekleme davranışı bir gün gerekirse ayrı bir fabrika ister — çünkü o
    zaman sınanan şey üyelik değil rol yetkisidir.
    """
    response = await client.post(
        f"/courses/{course_id}/members",
        json={"email": email, "role": "student"},
        headers=headers,
    )
    assert response.status_code == 201, response.text


async def create_topic(
    client: AsyncClient,
    headers: dict[str, str],
    course_id: UUID | str,
    name: str,
) -> UUID:
    """Derse konu ekler ve kimliğini döndürür. `headers` eğitmenin olmalıdır."""
    response = await client.post(
        f"/courses/{course_id}/topics", json={"name": name}, headers=headers
    )
    assert response.status_code == 201, response.text
    return UUID(response.json()["id"])


# ---------------------------------------------------------------------------
# Doğrudan SQL ile tohumlama
#
# Ingestion hattını çağırmak yerine satırlar doğrudan yazılır: retrieval, soru
# havuzu ve analitik testlerinde sınanan şey ayrıştırma değil, metnin üzerinde
# tam kontrol testin ne kanıtladığını okunur kılıyor. Yazan bağlantı çağırana
# bırakılır — chunks tablosunda `dou_app` için INSERT politikası bilinçli olarak
# YOKTUR, dolayısıyla çağrı yeri ya tablo sahibini ya da `dou_worker`'ı verir.
# ---------------------------------------------------------------------------


@dataclass(frozen=True, slots=True)
class SeededDocument:
    """Tohumlanan belge ve chunk'larının kimlikleri.

    İkisi birden döndürülüyor çünkü çağrı yerleri ikiye ayrılıyor: soru havuzu
    testleri chunk kimliklerini (soru bir chunk'a bağlanır), belge testleri ise
    belge kimliğini istiyor. Ayrı iki fabrika yazmak, aynı INSERT çiftini yine
    iki kez yazmak olurdu.
    """

    document_id: UUID
    chunk_ids: list[UUID]


_DOCUMENT_SQL = sql_text(
    "INSERT INTO documents (id, course_id, uploaded_by, file_name, file_type, "
    "storage_path, file_hash, byte_size, status, page_count, chunk_count) "
    "VALUES (:id, :course_id, :uploaded_by, :file_name, :file_type, :path, :hash, "
    "4096, 'completed', :pages, :chunks)"
)

# İki ayrı SQL, tek bir dinamik SQL değil: `embedding` sütunu `vector` tipine
# CAST istiyor ve gömüsüz tohumlamada bind edilecek bir vektör yok. Metni
# çalışma anında birleştirmek, sorguyu okunamaz ve denetlenemez kılardı.
_CHUNK_SQL = sql_text(
    "INSERT INTO chunks (id, course_id, document_id, chunk_index, page_number, "
    "section_title, content_type, text, token_count, embedding_space) "
    "VALUES (:id, :course_id, :document_id, :chunk_index, :page_number, "
    ":section_title, 'text', :text, :token_count, :embedding_space)"
)

_CHUNK_SQL_WITH_EMBEDDING = sql_text(
    "INSERT INTO chunks (id, course_id, document_id, chunk_index, page_number, "
    "section_title, content_type, text, token_count, embedding, embedding_space) "
    "VALUES (:id, :course_id, :document_id, :chunk_index, :page_number, "
    ":section_title, 'text', :text, :token_count, CAST(:embedding AS vector), "
    ":embedding_space)"
)


def _passage_rows(
    passages: Sequence[str | tuple[int | None, str]],
) -> list[tuple[int | None, str]]:
    """Pasajları `(sayfa, metin)` çiftlerine indirger.

    İki biçim de kabul ediliyor çünkü çağrı yerleri iki gruba ayrılıyor: sayfa
    numarasının kendisi iddianın parçası olduğu testler çifti açıkça yazar
    (atıf konumu oradan üretilir), sayfayı umursamayanlar düz metin verir ve
    sıradan türeyen `index + 1` yeter.
    """
    rows: list[tuple[int | None, str]] = []
    for index, item in enumerate(passages):
        if isinstance(item, tuple):
            rows.append(item)
        else:
            rows.append((index + 1, item))
    return rows


async def seed_document(
    engine: AsyncEngine,
    *,
    course_id: UUID,
    uploaded_by: UUID,
    passages: Sequence[str | tuple[int | None, str]],
    file_name: str = "isletim-sistemleri.pdf",
    embeddings: bool = False,
    embedding_space: str | None = None,
) -> SeededDocument:
    """Bir belge ve chunk'larını yazar.

    `embeddings` varsayılan olarak KAPALI: gömü üretmek yapılandırılmış
    sağlayıcıyı çağırır ve chunk'ları yalnız kaynak olsun diye yazan testler
    (soru havuzu, analitik) o maliyeti ödemek zorunda değil. Retrieval testleri
    açıkça açar — orada gömü sınanan şeyin ta kendisi.

    `embedding_space` varsayılan `None`: `0006` öncesi yazılmış satırların hâli.
    Bu, damgasız korpusun aramayı durdurmadığı iddiasını da sürekli sınayan
    bilinçli bir varsayılandır.

    `file_type` dosya adının uzantısından türetilir; ikisinin ayrı ayrı
    verilebilmesi yalnız tutarsız satır üretmeye yarardı.
    """
    rows = _passage_rows(passages)
    document_id = uuid4()
    chunk_ids = [uuid4() for _ in rows]
    _, _, extension = file_name.rpartition(".")
    file_type = extension or "pdf"

    if embeddings:
        from app.modules.ingestion.embedding import get_embedding_provider

        vectors = get_embedding_provider().embed_documents([body for _, body in rows])
    else:
        vectors = []

    async with engine.begin() as conn:
        await conn.execute(
            _DOCUMENT_SQL,
            {
                "id": document_id,
                "course_id": course_id,
                "uploaded_by": uploaded_by,
                "file_name": file_name,
                "file_type": file_type,
                "path": f"courses/{course_id}/{document_id.hex}.{file_type}",
                "hash": document_id.hex,
                "pages": len(rows),
                "chunks": len(rows),
            },
        )
        parameters = [
            {
                "id": chunk_id,
                "course_id": course_id,
                "document_id": document_id,
                "chunk_index": index,
                "page_number": page,
                "section_title": None,
                "text": body,
                "token_count": max(1, len(body.split())),
                "embedding_space": embedding_space,
            }
            for index, (chunk_id, (page, body)) in enumerate(zip(chunk_ids, rows, strict=True))
        ]
        if embeddings:
            for parameter, vector in zip(parameters, vectors, strict=True):
                parameter["embedding"] = str(vector)
            await conn.execute(_CHUNK_SQL_WITH_EMBEDDING, parameters)
        else:
            await conn.execute(_CHUNK_SQL, parameters)

    return SeededDocument(document_id=document_id, chunk_ids=chunk_ids)


_QUESTION_SQL = sql_text(
    "INSERT INTO questions (id, course_id, topic_id, type, payload, source_chunk_id, "
    "status, created_by, reviewed_by, reviewed_at) VALUES "
    "(:id, :course_id, :topic_id, CAST(:type AS question_type), CAST(:payload AS jsonb), "
    ":chunk_id, CAST(:status AS question_status), :created_by, :reviewed_by, "
    "CASE WHEN :status = 'draft' THEN NULL ELSE now() END)"
)


async def seed_question(
    engine: AsyncEngine,
    *,
    course_id: UUID,
    topic_id: UUID,
    source_chunk_id: UUID,
    payload: dict[str, Any],
    question_type: QuestionType = QuestionType.MCQ,
    status: str = "draft",
    created_by: UUID | None = None,
    reviewed_by: UUID | None = None,
) -> UUID:
    """Havuza soru yazar ve kimliğini döndürür.

    `reviewed_by` kontrolü burada, veritabanına gitmeden yapılıyor: aksi hâlde
    `questions_reviewed_consistency` CHECK'i tetiklenir ve test, kurulumunun
    eksik olduğunu okunmaz bir `IntegrityError` metninden anlamaya çalışır.
    """
    if status != "draft" and reviewed_by is None:
        raise AssertionError("questions_reviewed_consistency: reviewed_by verilmeli")

    question_id = uuid4()
    async with engine.begin() as conn:
        await conn.execute(
            _QUESTION_SQL,
            {
                "id": question_id,
                "course_id": course_id,
                "topic_id": topic_id,
                "type": question_type.value,
                "payload": json.dumps(payload),
                "chunk_id": source_chunk_id,
                "status": status,
                "created_by": created_by,
                "reviewed_by": reviewed_by,
            },
        )
    return question_id


# ---------------------------------------------------------------------------
# Yüklenebilir dosya baytları
#
# Sahte baytlarla çalışmak parser'ların sayfa/slayt numarasını doğru taşıdığını
# kanıtlamaz — sistemin tüm atıf vaadi buna dayanıyor. Bu yüzden gerçek dosyalar
# üretilir.
# ---------------------------------------------------------------------------


def make_pdf(pages: list[str]) -> bytes:
    """Her elemanı bir sayfa olan gerçek bir PDF üretir."""
    import pymupdf

    document = pymupdf.open()
    for body in pages:
        page = document.new_page()
        page.insert_text((72, 72), body, fontsize=11)
    data: bytes = document.tobytes()
    document.close()
    return data


def make_pptx(slides: list[tuple[str, str]]) -> bytes:
    """Her elemanı `(başlık, gövde)` olan gerçek bir PPTX üretir."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for title, body in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
        box.text_frame.text = body
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Sözleşme nesneleri
# ---------------------------------------------------------------------------


def make_chunk(
    *,
    text: str = "Deadlock için dört koşulun aynı anda sağlanması gerekir.",
    chunk_id: UUID | None = None,
    file_name: str = "isletim-sistemleri-hafta3.pdf",
    page_number: int | None = 7,
    section_title: str | None = "Deadlock Koşulları",
    dense_score: float = 0.82,
    fts_score: float = 0.4,
    fused_score: float = 0.032,
) -> RetrievedChunk:
    """Tek bir `RetrievedChunk` üretir.

    Dört dosyada beş ayrı adla duran sarmalayıcıların birleşimi. Her alan
    parametre çünkü her biri en az bir çağrı yerinde değişiyor: kanıt kapısı
    `dense_score`'a, izolasyon iddiaları `file_name`'e, atıf konumu
    `page_number`'a, şablon temizliği `section_title`'a bakıyor.

    Varsayılan `fused_score` GERÇEKÇİ bir RRF değeridir (~0.03) ve bu bilinçli:
    kanıt eşiği füzyon skoruna uygulansaydı her soru reddedilirdi, dolayısıyla
    şişirilmiş bir varsayılan o tuzağı testlerden gizlerdi.

    `chunk_id` verilmezse üretilir. Verilebilmesinin sebebi atıf doğrulaması:
    set-membership testleri "bu kimlik retrieve edilen kümede" ile "değil"
    arasındaki farkı kurabilmek için kimliği elde tutmak zorunda.

    `document_id` ve `slide_number` PARAMETRE DEĞİL: 48 çağrı yerinin hiçbiri
    onları vermiyordu. Slayt konumunu sınayan testler var ama hepsi gerçek
    chunk satırları üzerinden gidiyor, bu sahte üzerinden değil. Gerekirse
    o gün parametre olurlar.
    """
    return RetrievedChunk(
        chunk_id=chunk_id or uuid4(),
        document_id=uuid4(),
        file_name=file_name,
        page_number=page_number,
        slide_number=None,
        section_title=section_title,
        text=text,
        dense_score=dense_score,
        fts_score=fts_score,
        fused_score=fused_score,
    )


def sourced_answer(
    chunk: RetrievedChunk, text_value: str = "Dört koşul aynı anda sağlanır."
) -> GeneratedAnswer:
    """Verilen parçaya atıf yapan, cevaplanmış bir QA çıktısı.

    Atıf alanları modelin metninden değil parçanın metadata'sından üretilir
    (Anayasa I); sahte cevabın bu kuralı taklit etmesi, atıf zincirinin testte
    de gerçek girdiyle karşılaşmasını sağlar.
    """
    return GeneratedAnswer(
        status=AnswerStatus.ANSWERED,
        mode=ChatMode.QA,
        text=text_value,
        citations=[Citation(chunk.chunk_id, chunk.file_name, chunk.location, chunk.text[:60])],
    )


# ---------------------------------------------------------------------------
# Sahteler — `app/contracts.py` protokollerinin test uygulamaları
#
# Hepsi protokolleri uygular; sahte olmaları testlerin gerçek modüller indiğinde
# anlamını yitirmesine sebep olmaz (00_OKU_ONCE §2).
# ---------------------------------------------------------------------------


class FakeRetriever:
    """`contracts.Retriever` — hangi derse sorulursa sorulsun aynı parçaları döner.

    Ders bazlı ayrımı gerekmeyen testlerin sahtesi. İzolasyon iddiası kuran bir
    test bunu KULLANAMAZ: her zaman aynı sonucu veren bir arama, izolasyonu
    kanıtlayamaz. O senaryo için `CourseScopedRetriever` var.

    Çağrılar sayılır ve görülen sorgular kaydedilir; "kanıt yoksa modele hiç
    gidilmez" ve "arama denemeyle değil açılış sorusuyla yapılır" gibi iddialar
    ancak bu kayıtla sınanabilir.
    """

    def __init__(self, chunks: list[RetrievedChunk]) -> None:
        self._chunks = chunks
        self.calls = 0
        self.seen_course_ids: list[UUID] = []
        self.seen_queries: list[str] = []

    async def search(self, *, course_id: UUID, query: str, limit: int = 8) -> list[RetrievedChunk]:
        self.calls += 1
        self.seen_course_ids.append(course_id)
        self.seen_queries.append(query)
        return self._chunks[:limit]


class CourseScopedRetriever:
    """`contracts.Retriever` — kendisine hangi `course_id` ile gelindiğine BAKAR.

    İzolasyon iddiasının test edilebilmesi için gerekli: uç yetkilendirilmiş
    dersi geçirmezse arama boş döner ve test kırmızıya döner. Ders bilmeyen bir
    sahteyle aynı test hiçbir şey kanıtlamazdı.
    """

    def __init__(self) -> None:
        self.by_course: dict[UUID, list[RetrievedChunk]] = {}
        self.seen_course_ids: list[UUID] = []
        self.seen_queries: list[str] = []

    async def search(self, *, course_id: UUID, query: str, limit: int = 8) -> list[RetrievedChunk]:
        self.seen_course_ids.append(course_id)
        self.seen_queries.append(query)
        return self.by_course.get(course_id, [])[:limit]


class FakeGenerator:
    """`contracts.Generator` — her çağrıda sıradaki hazır cevabı döner.

    Liste tükenince sonuncusu tekrar edilir; "ısrarcı öğrenci" senaryoları aynı
    cevabı kaç tur alacağını önceden bilmek zorunda kalmasın diye.
    """

    def __init__(self, answers: list[GeneratedAnswer]) -> None:
        self._answers = answers
        self.calls = 0
        self.seen_stages: list[SocraticStage | None] = []
        #: Uç, öğrencinin denemesini gerçekten geçiriyor mu? Kaydedilmezse
        #: "ipucu denemeye göre şekilleniyor" iddiası test edilemez (Anayasa III).
        self.seen_attempts: list[str | None] = []

    async def generate(
        self,
        *,
        question: str,
        chunks: list[RetrievedChunk],
        mode: ChatMode,
        socratic_stage: SocraticStage | None = None,
        student_attempt: str | None = None,
    ) -> GeneratedAnswer:
        del question, chunks, mode
        self.calls += 1
        self.seen_stages.append(socratic_stage)
        self.seen_attempts.append(student_attempt)
        index = min(self.calls - 1, len(self._answers) - 1)
        answer = self._answers[index]
        # Üretim katmanı her çağrıda yeni nesne döndürür; zincir mutasyonu testler
        # arasında sızmasın diye kopyalanıyor.
        return GeneratedAnswer(
            status=answer.status,
            mode=answer.mode,
            text=answer.text,
            citations=list(answer.citations),
            socratic_stage=answer.socratic_stage,
        )


class FakeCitationGuardrail:
    """`contracts.Guardrail` — set-membership: her chunk_id retrieve edilen kümede olmalı.

    Adı `Fake` önekli çünkü `app.modules.guardrails.citation.CitationGuardrail`
    üretimde koşan gerçek halkadır ve testlerin bir kısmı ONU içe aktarıyor. İki
    ad aynı olsaydı, bir testin gerçek zinciri mi yoksa sahtesini mi sınadığı
    içe aktarma satırına bakmadan anlaşılamazdı.

    Kontrolün kendisinin doğruluğu gerçek halkanın testidir (T016); bu sahte,
    ondan bağımsız çalışması gereken şeritlerin bloklamaya UYDUĞUNU gösterir.
    """

    def check(self, answer: GeneratedAnswer, retrieved: list[RetrievedChunk]) -> GuardrailVerdict:
        allowed = {chunk.chunk_id for chunk in retrieved}
        dropped = [c.chunk_id for c in answer.citations if c.chunk_id not in allowed]
        remaining = [c for c in answer.citations if c.chunk_id in allowed]
        if not remaining:
            return GuardrailVerdict(
                blocked=True, reason="citation:no_valid_citation", dropped_citations=dropped
            )
        return GuardrailVerdict(blocked=False, dropped_citations=dropped)


class Pipeline:
    """Sohbet ucuna takılacak hat: ders bazlı sahte arama + sahte üretim + atıf kontrolü.

    Atıf kontrolü sahte ama GERÇEK set-membership'i yapıyor: kaynaksız cevabın
    kullanıcıya gitmediği iddiası ürünün tezi ve o iddianın sohbet ucu
    testlerinde de gösterilebilmesi gerekiyor.
    """

    def __init__(self) -> None:
        self.retriever = CourseScopedRetriever()
        self.generator = FakeGenerator([])

    def serve(self, course_id: UUID | str, chunk: RetrievedChunk) -> RetrievedChunk:
        """Verilen derse bir parça ekler ve onu geri verir.

        Parçayı geri vermesi kolaylık değil: çağıran genellikle hemen ardından
        o parçaya atıf yapan bir cevap kuruyor ve kimliğin elde kalması gerekiyor.
        """
        self.retriever.by_course.setdefault(UUID(str(course_id)), []).append(chunk)
        return chunk

    def answers(self, *answers: GeneratedAnswer) -> None:
        """Hazır cevapları bağlar ve hattı yeniden takar.

        Yeniden takmak zorunlu: `set_pipeline` üreteci DEĞER olarak alıyor,
        dolayısıyla yeni bir `FakeGenerator` kurulduğunda uç eskisini tutmaya
        devam ederdi.
        """
        self.generator = FakeGenerator(list(answers))
        install_pipeline(self)


def install_pipeline(pipeline: Pipeline) -> None:
    """Hattı sohbet ucuna takar.

    İçe aktarma gövdede: `app.api.chat` ayarları ve veritabanı katmanını
    çekiyor, bu modül ise ondan bağımsız kullanılabilmeli (parça ve tohumlama
    fabrikaları uca hiç dokunmuyor).
    """
    from app.api.chat import set_pipeline

    set_pipeline(
        retriever_factory=lambda _session: pipeline.retriever,
        generator=pipeline.generator,
        guardrails=[FakeCitationGuardrail()],
    )


# ---------------------------------------------------------------------------
# Soru havuzu korpusu ve payload üreteçleri
#
# Havuz (`test_assessment`), sınav (`test_exams`), blueprint (`test_blueprint`)
# ve hız sınırı (`test_rate_limit`) testleri aynı korpusu görmek zorunda: sınav
# testleri havuzun ürettiği soruların aynısını okur, ayrı bir korpus iki
# dosyanın sessizce ayrışmasına açık kapı bırakırdı. Eşik (üç dosya) aşıldı.
# ---------------------------------------------------------------------------


DEADLOCK_TEXTS = [
    "Deadlock için dört koşulun aynı anda sağlanması gerekir: karşılıklı dışlama, "
    "tut ve bekle, kesilemezlik ve döngüsel bekleme.",
    "Kesme (preemption) zorunluluğu bir deadlock koşulu değildir; tam tersine "
    "kaynağın zorla geri alınabilmesi deadlock'u kırar.",
]


def mcq_payload(chunk_ids: list[UUID]) -> dict[str, Any]:
    """Dört şıklı örnek MCQ. Çeldirici kaynakları gerçek chunk'lara bağlıdır."""
    return {
        "stem": "Deadlock'un oluşması için aşağıdakilerden hangisi gerekli DEĞİLDİR?",
        "options": [
            {"key": "A", "text": "Karşılıklı dışlama"},
            {"key": "B", "text": "Döngüsel bekleme"},
            {"key": "C", "text": "Kesme zorunluluğu"},
            {"key": "D", "text": "Tut ve bekle"},
        ],
        "answer_key": "C",
        "distractor_sources": {
            "A": str(chunk_ids[0]),
            "B": str(chunk_ids[0]),
            "D": str(chunk_ids[1 % len(chunk_ids)]),
        },
        "explanation": "Dört koşul arasında kesme zorunluluğu yoktur.",
    }


def short_answer_payload() -> dict[str, Any]:
    return {
        "prompt": "Deadlock'un dört koşulundan döngüsel beklemenin adı nedir?",
        "answer_key": "döngüsel bekleme",
        "format": AnswerFormat.SHORT_ANSWER.value,
        "accepted_answers": ["döngüsel bekleme", "circular wait"],
    }


ESSAY_PAYLOAD = {
    "prompt": "Deadlock'un dört koşulunu açıklayın.",
    "answer_key": "Karşılıklı dışlama, tut ve bekle, kesilemezlik, döngüsel bekleme.",
    "key_points": ["karşılıklı dışlama", "tut ve bekle", "döngüsel bekleme"],
    "rubric": [{"point": "Dört koşulu sayar", "weight": 100}],
}


def _mcq_response(source_chunk_id: UUID | str, *, count: int = 1) -> str:
    return json.dumps(
        {
            "questions": [
                {
                    "stem": f"Deadlock koşullarından hangisi yanlıştır? ({index})",
                    "options": [
                        {"key": "A", "text": "Karşılıklı dışlama"},
                        {"key": "B", "text": "Döngüsel bekleme"},
                        {"key": "C", "text": "Kesme zorunluluğu"},
                        {"key": "D", "text": "Tut ve bekle"},
                    ],
                    "answer_key": "C",
                    "explanation": "Dört koşul arasında kesme zorunluluğu yoktur.",
                    "source_chunk_id": str(source_chunk_id),
                }
                for index in range(count)
            ]
        }
    )


def retrieved(chunk_ids: list[UUID], texts: list[str]) -> list[RetrievedChunk]:
    """Soru üretimine verilecek arama sonucu.

    `dense_score`/`fts_score` AÇIKÇA sıfırlanıyor. `make_chunk`'ın varsayılanları
    sohbet tarafının gerçekçi değerleridir (dense 0.82) ve kanıt eşiği oraya
    bakar; soru üretimi ise eşiğe hiç bakmaz, sıralamayı yalnız `fused_score`
    belirler. Varsayılanlar devralınsaydı bu testler ölçmedikleri bir eşiğe
    bağlanır ve eşik bir gün değiştiğinde sebepsiz kırılırlardı.
    """
    return [
        make_chunk(
            chunk_id=chunk_id,
            text=body,
            file_name="isletim-sistemleri.pdf",
            page_number=index + 1,
            section_title=None,
            dense_score=0.0,
            fts_score=0.0,
            fused_score=1.0 - index * 0.1,
        )
        for index, (chunk_id, body) in enumerate(zip(chunk_ids, texts, strict=True))
    ]


# ---------------------------------------------------------------------------
# Sahte LLM sağlayıcıları
# ---------------------------------------------------------------------------


class FakeCompletion:
    """Sırayla hazır yanıt döndüren sahte LLM. Çağrı sayısı sayılır."""

    def __init__(self, *responses: str) -> None:
        self._responses = list(responses)
        self.calls = 0

    async def complete(self, *, system: str, user: str) -> str:
        del system, user
        self.calls += 1
        index = min(self.calls - 1, len(self._responses) - 1)
        return self._responses[index]


class ScriptedLlm:
    """Sırayla verilen ham metinleri döndüren sahte sağlayıcı.

    `app.modules.generation.fake.FakeLlmClient` gerçekçi cevap üretir; bu ise
    BELİRLİ bir bozuk/kötü çıktıyı zincire sokmak için var. İkisi ayrı çünkü
    "model şunu döndürürse ne olur" sorusu, gerçekçi bir sağlayıcıyla sorulamaz.

    İki dosyada kopyaydı (`test_guardrails`, `test_generation`) ve kopyalar
    ayrışmıştı: biri istekleri kaydediyordu, diğeri kaydetmiyordu. En genel
    biçim alındı (Anayasa XI istisna b): `requests` kaydı kalan kullanıcı için
    yalnızca boş bir liste taşır.
    """

    def __init__(self, *payloads: str) -> None:
        self._payloads = payloads or ("",)
        self.calls = 0
        self.requests: list[object] = []

    async def complete(self, request: object) -> object:
        from app.modules.generation.llm import LlmCompletion

        index = min(self.calls, len(self._payloads) - 1)
        self.calls += 1
        self.requests.append(request)
        return LlmCompletion(text=self._payloads[index], provider="scripted", model="scripted/test")


class BlockingGenerator:
    """`contracts.Generator` — ilk çağrıda `entered`'ı işaretler, `release`'e dek bekler.

    Yarış pencerelerini ZORLA açmak için var: üretim sürerken ikinci bir isteğin
    araya girebildiği anı kurar; sınav kilidi (`test_exam_lock`) ve eşzamanlılık
    kapısı (`test_role_aware_agent_application_guards`) bunun üstünde iddia kurar.

    İki dosyada gömülü kopyaydı ve kopyalar ayrışmıştı: biri çağrı sayıp yalnız
    ilkinde bekliyordu, diğeri saymıyordu ama tek çağrı aldığı için fark
    gözlemlenemezdi. En genel biçim alındı (Anayasa XI istisna b): sayaç var,
    yalnız ilk çağrı bekler.
    """

    def __init__(
        self, chunk: RetrievedChunk, *, entered: asyncio.Event, release: asyncio.Event
    ) -> None:
        self._chunk = chunk
        self._entered = entered
        self._release = release
        self.calls = 0

    async def generate(
        self,
        *,
        question: str,
        chunks: list[RetrievedChunk],
        mode: ChatMode,
        socratic_stage: SocraticStage | None = None,
        student_attempt: str | None = None,
    ) -> GeneratedAnswer:
        del question, chunks, mode, socratic_stage, student_attempt
        self.calls += 1
        if self.calls == 1:
            self._entered.set()
            await self._release.wait()
        return sourced_answer(self._chunk)


# ---------------------------------------------------------------------------
# Sınav kurulumu — ders + roller + onaylı havuz + oturum yardımcıları
#
# Uç sarmalayıcılarından yalnız `start` burada: oturum açmayı dört dosya
# kullanıyor. `answer`/`finish`/`hint` ise yalnız `test_exams`'e ait ve orada
# kaldılar — o üç uç durum kodu iddia etmeden ham `Response` döndürmek zorunda.
# ---------------------------------------------------------------------------


EXAM_DURATION_SECONDS = 20 * 60  # Settings.exam_duration_minutes varsayılanı


class ExamFixture:
    """Bir ders, bir eğitmen, bir öğrenci ve onaylı bir soru havuzu.

    `course_id` `str` KALIYOR, `UUID` değil: `tests/test_exam_lock.py` bu sınıfı
    içe aktarıyor ve kendi yardımcılarını `str` olarak imzalıyor. Fabrika artık
    `UUID` döndürdüğü için dönüşüm `build_course`'un içinde yapılıyor — tipi
    burada değiştirmek, bu turda dokunulmayan imzaları yalancı çıkarırdı.
    """

    def __init__(
        self,
        *,
        course_id: str,
        instructor: dict[str, str],
        instructor_id: UUID,
        student: dict[str, str],
        student_id: UUID,
        topic_id: UUID,
        chunk_ids: list[UUID],
        question_ids: list[UUID],
    ) -> None:
        self.course_id = course_id
        self.instructor = instructor
        self.instructor_id = instructor_id
        self.student = student
        self.student_id = student_id
        self.topic_id = topic_id
        self.chunk_ids = chunk_ids
        self.question_ids = question_ids


async def build_course(
    client: AsyncClient,
    users: UserFactory,
    admin_engine: AsyncEngine,
    *,
    approved: int = 2,
    drafts: int = 0,
    short_answer: bool = False,
) -> ExamFixture:
    instructor_id = await users.create("ayse@dogus.edu.tr")
    instructor = users.auth(instructor_id)
    student_id = await users.create("burak@dogus.edu.tr")
    course_id = await create_course(client, instructor, "COME301")
    await enroll_student(client, instructor, course_id, "burak@dogus.edu.tr")
    topic_id = await create_topic(client, instructor, course_id, "Deadlock")
    seeded = await seed_document(
        admin_engine, course_id=course_id, uploaded_by=instructor_id, passages=DEADLOCK_TEXTS
    )
    chunk_ids = seeded.chunk_ids

    question_ids: list[UUID] = []
    for index in range(approved):
        use_short = short_answer and index == 0
        question_ids.append(
            await seed_question(
                admin_engine,
                course_id=course_id,
                topic_id=topic_id,
                source_chunk_id=chunk_ids[index % len(chunk_ids)],
                payload=short_answer_payload() if use_short else mcq_payload(chunk_ids),
                question_type=QuestionType.OPEN if use_short else QuestionType.MCQ,
                status="approved",
                reviewed_by=instructor_id,
            )
        )
    for _ in range(drafts):
        await seed_question(
            admin_engine,
            course_id=course_id,
            topic_id=topic_id,
            source_chunk_id=chunk_ids[0],
            payload=mcq_payload(chunk_ids),
        )

    return ExamFixture(
        course_id=str(course_id),
        instructor=instructor,
        instructor_id=instructor_id,
        student=users.auth(student_id),
        student_id=student_id,
        topic_id=topic_id,
        chunk_ids=chunk_ids,
        question_ids=question_ids,
    )


async def start(client: AsyncClient, fixture: ExamFixture, mode: str) -> dict:
    response = await client.post(
        f"/courses/{fixture.course_id}/exams", json={"mode": mode}, headers=fixture.student
    )
    assert response.status_code == 201, response.text
    return response.json()


async def rewind(admin_engine: AsyncEngine, session_id: str, *, minutes: int) -> None:
    """Oturumu geriye alır: süre gerçekten geçmiş gibi davranır."""
    async with admin_engine.begin() as conn:
        await conn.execute(
            sql_text(
                "UPDATE exam_sessions SET started_at = started_at - make_interval(mins => :m), "
                "expires_at = expires_at - make_interval(mins => :m) WHERE id = :id"
            ),
            {"m": minutes, "id": UUID(session_id)},
        )


# ---------------------------------------------------------------------------
# Yavaş gömü sağlayıcısı — event loop ve sağlık yoklaması testleri
# ---------------------------------------------------------------------------


#: Sahte sağlayıcının her çağrıda harcadığı süre.
SLOW_CALL_SECONDS = 0.4


class SlowSyncEmbeddingProvider(HashingEmbeddingProvider):
    """Yavaş ama GIL'i BIRAKAN senkron sağlayıcı.

    `time.sleep` bilinçli bir seçim: GIL'i bırakır, yani ONNX çıkarımının ve
    PyMuPDF ayrıştırmasının dürüst vekilidir. Saf Python döngüsüyle beklenseydi
    GIL bırakılmaz, `to_thread` sarması yerindeyken de test kırmızı yanar ve
    kusur yanlış yerde aranırdı.

    `HashingEmbeddingProvider`'dan türüyor ki `vector_space.space_of` onu
    tanısın ve damga normal test sağlayıcısıyla aynı kalsın; farklı bir uzay
    kimliği üretseydi retrieval kapısı ölçümden önce devreye girerdi.
    """

    def __init__(self, delay: float = SLOW_CALL_SECONDS) -> None:
        super().__init__()
        self._delay = delay
        self.calls = 0

    def embed_documents(self, texts: Sequence[str]) -> list[list[float]]:
        self.calls += 1
        time.sleep(self._delay)
        return super().embed_documents(texts)

    def embed_query(self, text: str) -> list[float]:
        self.calls += 1
        time.sleep(self._delay)
        return super().embed_query(text)
