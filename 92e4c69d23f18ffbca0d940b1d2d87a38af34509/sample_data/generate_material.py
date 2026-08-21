#!/usr/bin/env python3
"""Ders materyalini kaynak Markdown'dan PDF ve PPTX'e çevirir (R2, İş 0).

Neden bir betik: `sample_data/isletim-sistemleri` v1 paketindeki PDF'ler pandoc +
XeLaTeX ile üretilmişti ve o araç zinciri bu makinede yok. Materyali elle üretilmiş
ikili dosyalar olarak commit etmek, "bu sayfa neden burada bitiyor" sorusunun cevabını
kaybettirir. Kaynak Markdown depoda durur, ikili dosya ondan **yeniden üretilebilir.**

Sayfa sınırı ölçüm için belirleyicidir: chunking bir chunk'ın iki sayfayı
birleştirmesine izin vermez (ARCHITECTURE §3) ve gold set'in kalıcı kimliği
`(dosya, sayfa)` çiftidir. Bu yüzden sayfa bölmesi **otomatik akışa bırakılmaz**,
kaynakta `<!-- sayfa -->` ile açıkça yazılır. Akışa bırakılsaydı fontun birkaç
piksellik farkı bir cümleyi sonraki sayfaya atar ve gold set sessizce koparadı.

Kullanım (depo kökünden):

    cd apps/api && uv run python ../../sample_data/generate_material.py
    cd apps/api && uv run python ../../sample_data/generate_material.py --check

`--check` hiçbir dosya yazmaz; üretilen çıktı depodakiyle aynı mı diye bakar.

**Font:** metin bizim, font değil. Türkçe `ı/İ/ş/ğ` glifleri Base-14 fontlarda yok;
işletim sisteminin Unicode fontu kullanılır ve PDF'e yalnız alt kümesi gömülür — bir
belgeyi yazdırırken olan şeyin aynısı. Font dosyası depoya konmaz.
"""

from __future__ import annotations

import argparse
import html
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

MATERIAL_DIR = Path(__file__).resolve().parent / "isletim-sistemleri"

#: Türkçe glifleri olan aday fontlar, tercih sırasıyla. İlk bulunan kullanılır.
#: Arial Unicode MS bilinçli olarak listede DEĞİL: 22 MB'lık CJK kapsamıyla geliyor ve
#: alt kümeleme öncesi 15 MB'lık PDF üretiyordu. Latin Genişletilmiş-A yeten bir font
#: hem küçük hem yeterli.
FONT_CANDIDATES = (
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
)

#: Fontta bulunması ZORUNLU gliflerin denetimi. Türkçe noktalı/noktasız i ayrımı
#: (Anayasa V) ve materyalde geçen tipografik işaretler eksikse PDF sessizce boş kutu
#: basar — ekranda değil, ancak ingest edilen metne bakınca fark edilir.
REQUIRED_GLYPHS = "ıİşŞğĞçÇöÖüÜâÂ—•→"

PAGE_BREAK = re.compile(r"^<!--\s*sayfa\s*-->\s*$", re.MULTILINE)
SLIDE_BREAK = re.compile(r"^<!--\s*slayt\s*-->\s*$", re.MULTILINE)
FRONT_MATTER = re.compile(r"\A---\n(.*?)\n---\n", re.DOTALL)

PAGE_CSS = """
* { font-family: gövde; }
body { font-size: 9.5pt; line-height: 1.32; }
h1 { font-size: 13pt; margin-top: 0; margin-bottom: 4pt; }
h2 { font-size: 11pt; margin-top: 7pt; margin-bottom: 3pt; }
h3 { font-size: 10pt; margin-top: 6pt; margin-bottom: 2pt; }
p  { margin-top: 0; margin-bottom: 5pt; }
li { margin-bottom: 2pt; }
pre { font-size: 8.5pt; margin-top: 3pt; margin-bottom: 5pt; }
table { font-size: 9pt; }
td, th { padding: 2pt 5pt 2pt 0; text-align: left; }
.baslik { font-size: 15pt; margin-bottom: 1pt; }
.altbaslik { font-size: 10pt; margin-bottom: 9pt; }
"""


class RenderError(RuntimeError):
    """Materyal üretilemedi — sessizce bozuk dosya yazmaktansa dur."""


def find_font() -> str:
    import pymupdf

    for candidate in FONT_CANDIDATES:
        if not Path(candidate).exists():
            continue
        font = pymupdf.Font(fontfile=candidate)
        missing = [glyph for glyph in REQUIRED_GLYPHS if not font.has_glyph(ord(glyph))]
        if missing:
            print(f"  atlandı {candidate}: eksik glif {''.join(missing)}", file=sys.stderr)
            continue
        return candidate
    raise RenderError(
        "Türkçe gliflere sahip bir font bulunamadı. Denenen yollar:\n  "
        + "\n  ".join(FONT_CANDIDATES)
        + "\nFONT_CANDIDATES listesine sistemdeki bir Unicode TTF ekleyin."
    )


@dataclass(slots=True)
class Source:
    """Kaynak Markdown: başlık bilgisi + açıkça bölünmüş parçalar."""

    path: Path
    title: str
    subtitle: str
    kind: str  # "pdf" | "pptx"
    parts: list[str] = field(default_factory=list)


def parse_front_matter(text: str) -> tuple[dict[str, str], str]:
    match = FRONT_MATTER.match(text)
    if not match:
        return {}, text
    meta: dict[str, str] = {}
    for line in match.group(1).splitlines():
        key, _, value = line.partition(":")
        if key.strip():
            meta[key.strip()] = value.strip().strip('"')
    return meta, text[match.end() :]


def load_source(path: Path) -> Source | None:
    """Kaynağı okur. `format:` beyan etmeyen dosya bu betiğin işi DEĞİLDİR.

    v1 paketindeki beş Markdown (pandoc/XeLaTeX ile üretilmiş PDF'lerin kaynağı)
    sayfa bölmesi taşımıyor. Bu betik onları da işleseydi her biri tek sayfalık bir
    PDF'e dönüşür ve `holdout.json`'daki 91 kaynak referansının sayfa numaraları
    sessizce koparadı. Beyan zorunlu tutuluyor: karışıklık pahalı, beyan ucuz.
    """
    meta, body = parse_front_matter(path.read_text(encoding="utf-8"))
    kind = meta.get("format")
    if kind is None:
        return None
    if kind not in {"pdf", "pptx"}:
        raise RenderError(f"{path.name}: 'format' pdf ya da pptx olmalı, '{kind}' geldi.")
    splitter = SLIDE_BREAK if kind == "pptx" else PAGE_BREAK
    parts = [part.strip() for part in splitter.split(body)]
    parts = [part for part in parts if part]
    if not parts:
        raise RenderError(f"{path.name}: içerik yok.")
    return Source(
        path=path,
        title=meta.get("title", path.stem),
        subtitle=meta.get("subtitle", ""),
        kind=kind,
        parts=parts,
    )


# ---------------------------------------------------------------------------
# Markdown -> HTML (bilinçli olarak dar bir alt küme)
# ---------------------------------------------------------------------------
#
# Tam bir Markdown uygulaması gerekmiyor ve bağımlılık eklemeye değmez; kaynak
# dosyaları bu betikle birlikte yazıldı, desteklenmeyen bir imle karşılaşırsa betik
# susmak yerine hata verir (aşağıdaki tablo/kod dalları dışında her satır paragraf).

_INLINE_CODE = re.compile(r"`([^`]+)`")
_BOLD = re.compile(r"\*\*([^*]+)\*\*")
_ITALIC = re.compile(r"(?<!\*)\*([^*]+)\*(?!\*)")


def inline(text: str) -> str:
    escaped = html.escape(text)
    escaped = _INLINE_CODE.sub(r"<b>\1</b>", escaped)
    escaped = _BOLD.sub(r"<b>\1</b>", escaped)
    return _ITALIC.sub(r"<i>\1</i>", escaped)


def markdown_to_html(text: str) -> str:
    out: list[str] = []
    lines = text.splitlines()
    index = 0
    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if not stripped:
            index += 1
            continue

        if stripped.startswith("```"):
            block: list[str] = []
            index += 1
            while index < len(lines) and not lines[index].strip().startswith("```"):
                block.append(html.escape(lines[index]))
                index += 1
            index += 1
            out.append("<pre>" + "\n".join(block) + "</pre>")
            continue

        if stripped.startswith("|"):
            rows: list[str] = []
            while index < len(lines) and lines[index].strip().startswith("|"):
                rows.append(lines[index].strip())
                index += 1
            out.append(render_table(rows))
            continue

        if stripped.startswith("### "):
            out.append(f"<h3>{inline(stripped[4:])}</h3>")
            index += 1
            continue
        if stripped.startswith("## "):
            out.append(f"<h2>{inline(stripped[3:])}</h2>")
            index += 1
            continue
        if stripped.startswith("# "):
            out.append(f"<h1>{inline(stripped[2:])}</h1>")
            index += 1
            continue

        if stripped.startswith(("- ", "* ")):
            items: list[str] = []
            while index < len(lines) and lines[index].strip().startswith(("- ", "* ")):
                item = lines[index].strip()[2:]
                index += 1
                # Girintili devam satırları aynı maddeye ait.
                while index < len(lines) and lines[index].startswith("  ") and lines[index].strip():
                    item += " " + lines[index].strip()
                    index += 1
                items.append(f"<li>{inline(item)}</li>")
            out.append("<ul>" + "".join(items) + "</ul>")
            continue

        paragraph = [stripped]
        index += 1
        while index < len(lines) and lines[index].strip() and not _starts_block(lines[index]):
            paragraph.append(lines[index].strip())
            index += 1
        out.append(f"<p>{inline(' '.join(paragraph))}</p>")

    return "".join(out)


def _starts_block(line: str) -> bool:
    stripped = line.strip()
    return stripped.startswith(("#", "- ", "* ", "|", "```"))


def render_table(rows: list[str]) -> str:
    def cells(row: str) -> list[str]:
        return [cell.strip() for cell in row.strip().strip("|").split("|")]

    body = [row for row in rows if not set(row.replace("|", "").strip()) <= {"-", ":", " "}]
    if not body:
        return ""
    head, rest = body[0], body[1:]
    out = ["<table><tr>"]
    out += [f"<th>{inline(cell)}</th>" for cell in cells(head)]
    out.append("</tr>")
    for row in rest:
        out.append("<tr>" + "".join(f"<td>{inline(cell)}</td>" for cell in cells(row)) + "</tr>")
    out.append("</table>")
    return "".join(out)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

A4 = (595, 842)
MARGIN = 54


def render_pdf(source: Source, font_path: str) -> bytes:
    import io

    import pymupdf

    archive = pymupdf.Archive()
    archive.add(Path(font_path).read_bytes(), "govde.ttf")
    css = "@font-face {font-family: gövde; src: url(govde.ttf);}" + PAGE_CSS

    mediabox = pymupdf.Rect(0, 0, *A4)
    frame = pymupdf.Rect(MARGIN, MARGIN, A4[0] - MARGIN, A4[1] - MARGIN)
    buffer = io.BytesIO()
    writer = pymupdf.DocumentWriter(buffer)

    for number, part in enumerate(source.parts, start=1):
        body = markdown_to_html(part)
        if number == 1:
            body = (
                f'<p class="baslik"><b>{html.escape(source.title)}</b></p>'
                f'<p class="altbaslik">{html.escape(source.subtitle)}</p>' + body
            )
        story = pymupdf.Story(html=f"<body>{body}</body>", user_css=css, archive=archive)
        device = writer.begin_page(mediabox)
        more, _ = story.place(frame)
        story.draw(device)
        writer.end_page()
        if more:
            # Sayfa taşarsa sessizce kırpılır ve kaynak Markdown ile PDF ayrışır:
            # gold set'in sayfa numarası yanlış olur. Fail-closed (Anayasa IV).
            raise RenderError(
                f"{source.path.name} sayfa {number}: içerik sayfaya sığmadı. "
                "Kaynakta bir '<!-- sayfa -->' bölmesi ekleyin ya da metni kısaltın."
            )
    writer.close()

    document = pymupdf.open("pdf", buffer.getvalue())
    document.set_metadata(
        {
            "title": source.title,
            "author": "DOU-Synapse takımı",
            "subject": source.subtitle,
            "creator": "sample_data/generate_material.py (PyMuPDF Story)",
            "producer": "PyMuPDF",
            "keywords": "işletim sistemleri, ders materyali, örnek",
        }
    )
    # Font alt kümeleme olmadan tüm TTF gömülüyor ve 3 sayfalık bir ders notu megabaytlar
    # tutuyor. Depoya girecek dosyanın boyutu, yükleme ucunun 20 MB sınırıyla da ilgili.
    document.subset_fonts(verbose=False)
    data: bytes = document.tobytes(garbage=4, deflate=True)
    document.close()
    return data


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def render_pptx(source: Source) -> bytes:
    import io

    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    opening = presentation.slides.add_slide(title_layout)
    opening.shapes.title.text = source.title
    opening.placeholders[1].text = f"{source.subtitle}\nDOU-Synapse örnek ders materyali"

    for part in source.parts:
        lines = [line.rstrip() for line in part.splitlines() if line.strip()]
        if not lines:
            continue
        heading = lines[0].lstrip("#").strip()
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = heading
        frame = slide.placeholders[1].text_frame
        frame.word_wrap = True

        first = True
        for line in lines[1:]:
            text = line.strip()
            level = 0
            if text.startswith(("- ", "* ")):
                text = text[2:]
            elif line.startswith("  ") and text.startswith(("- ", "* ")):
                text, level = text[2:], 1
            text = _BOLD.sub(r"\1", _INLINE_CODE.sub(r"\1", text))
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            paragraph.text = text
            paragraph.level = level
            for run in paragraph.runs:
                run.font.size = Pt(15)
            first = False

    # v1 paketindeki sunum python-pptx'in varsayılan şablon tarihini (2013-01-27)
    # taşıyordu ve PR incelemesinde "kopyalanmış şablon mu" sorusunu doğurdu
    # (docs/team/PR_INCELEME_2026-08-06.md). Üretimi kim yaptıysa dosyada yazsın.
    core = presentation.core_properties
    core.title = source.title
    core.author = "DOU-Synapse takımı"
    core.comments = "sample_data/generate_material.py ile üretildi"
    core.subject = source.subtitle

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Ders materyalini kaynaktan üretir")
    parser.add_argument("--material", type=Path, default=MATERIAL_DIR)
    parser.add_argument("--only", help="Yalnız adı bu metni içeren kaynakları üret.")
    parser.add_argument(
        "--check", action="store_true", help="Yazma; depodaki çıktıyla aynı mı diye bak."
    )
    args = parser.parse_args(argv)

    font_path = find_font()
    sources = sorted(args.material.glob("*.md"))
    if args.only:
        sources = [path for path in sources if args.only in path.name]

    generated = 0
    skipped = 0
    stale: list[str] = []
    for path in sources:
        source = load_source(path)
        if source is None:
            skipped += 1
            continue
        if source.kind == "pdf":
            data, suffix = render_pdf(source, font_path), ".pdf"
        else:
            data, suffix = render_pptx(source), ".pptx"
        target = path.with_suffix(suffix)

        if args.check:
            # PDF/PPTX ikili çıktısı bit düzeyinde kararlı değil (zaman damgası, zip
            # sırası). Karşılaştırma parça SAYISI üzerinden: sayfa/slayt sınırı gold
            # set'i bağlayan şeydir, dosyanın baytları değil.
            if not target.exists():
                stale.append(f"{target.name}: yok")
            else:
                expected = len(source.parts) + (1 if source.kind == "pptx" else 0)
                actual = count_parts(target)
                if actual != expected:
                    stale.append(f"{target.name}: {actual} parça, kaynak {expected} diyor")
            continue

        target.write_bytes(data)
        generated += 1
        print(f"  {target.name}: {len(source.parts)} parça, {len(data) // 1024} KB")

    if args.check:
        for line in stale:
            print(f"  BAYAT {line}", file=sys.stderr)
        print(f"{len(sources) - skipped} kaynak denetlendi, {len(stale)} bayat, {skipped} atlandı.")
        return 1 if stale else 0

    print(f"{generated} dosya üretildi, {skipped} kaynak atlandı (v1: 'format' beyanı yok).")
    return 0


def count_parts(path: Path) -> int:
    if path.suffix == ".pdf":
        import pymupdf

        with pymupdf.open(path) as document:
            return int(document.page_count)
    from pptx import Presentation

    return len(Presentation(str(path)).slides)


if __name__ == "__main__":
    sys.exit(main())
